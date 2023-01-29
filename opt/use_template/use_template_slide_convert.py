import pptx
import collections.abc
import os

def use_template_slide_convert(target_path, file_name, save_folder):
    prs = pptx.Presentation(target_path)
    template_pptx = pptx.Presentation('template_slide.pptx')
    title_slide_layout = template_pptx.slide_masters[0].slide_layouts[0]
    lyric_slide_layout = template_pptx.slide_masters[0].slide_layouts[1]
    new_slide = template_pptx.slides.add_slide(title_slide_layout)
    has_picture = False

    # Title Slide
    for index, shape in enumerate(prs.slides[0].shapes):
        if type(shape)==pptx.shapes.picture.Picture or type(shape)==pptx.shapes.placeholder.PlaceholderPicture:
            has_picture = True
            break
        has_composition = '作詞' in shape.text_frame.text
        if index == 0:
            text = shape.text_frame.text
        if index > 0 and has_composition:
            new_slide.shapes[1].text_frame.text = shape.text_frame.text
        if index > 0 and not(has_composition):
            text = '\n'.join([text,shape.text_frame.text])
        new_slide.shapes[0].text_frame.text = text

    # Lyric slides
    # First version | for shape in slide.shapes:
    # Can't get using [1:]
    slides_total_num = len(prs.slides)
    for index in range(1,slides_total_num):
        if has_picture:
            with open('has_picture_slide.txt', 'a') as f:
                print(file_name, file=f)
            break
        # Get each target pptx slides
        slide = prs.slides[index]
        # Add blank Slide
        new_slide = template_pptx.slides.add_slide(lyric_slide_layout)
        # Check slide has text_frame
        for shape_index, shape in enumerate(slide.shapes):
            # Stop when slide has picture
            if type(shape)==pptx.shapes.picture.Picture or type(shape)==pptx.shapes.placeholder.PlaceholderPicture:
                has_picture = True
                break
            if not shape.has_text_frame:
                continue
            # Get TextFrame from various Shape objects
            # strip() is to delete pointless line breaks
            if shape_index == 0:
                text = shape.text_frame.text.strip()
            # If each row is created in a separate object
            if shape_index > 0:
                text = '\n'.join([text,shape.text_frame.text.strip()])
        new_slide.shapes[0].text_frame.text = text
    if has_picture:
        return True
    else:
        template_pptx.save(save_folder + '/' + file_name)
        return False