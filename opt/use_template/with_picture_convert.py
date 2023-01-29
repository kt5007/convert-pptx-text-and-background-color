import pptx
import collections.abc
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

def with_picture_convert(target_file_path,save_file_path):
    prs = pptx.Presentation(target_file_path)
    slides = prs.slides
    template_pptx = pptx.Presentation('template_slide.pptx')
    title_slide_layout = template_pptx.slide_masters[0].slide_layouts[0]
    lyric_slide_layout = template_pptx.slide_masters[0].slide_layouts[1]
    for slide in slides:
        # Change background color to white
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255,255,255)
        shape_positions = []
        for shape in slide.shapes:
            if not shape.has_text_frame: 
                continue
            # print(shape.text_frame.paragraphs[0].font.size)
            # Get attributes and text from existing shape
            x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
            shape_positions.append({"x":x, "y":y, "cx":cx, "cy":cy, "text":shape.text_frame.text})
            # Remove existing shape
            sp = shape._sp             # --- get reference to XML element for shape ---
            sp.getparent().remove(sp)  # --- remove that shape element from its tree ---
        for shape_position in shape_positions:
            new_shape = slide.shapes.add_textbox(shape_position["x"], shape_position["y"] ,shape_position["cx"], shape_position["cy"])
            new_shape.text_frame.text = shape_position["text"]
            for paragraph in new_shape.text_frame.paragraphs:
                paragraph.font.color.rgb = RGBColor(0, 0, 0)
                paragraph.font.size = Pt(42)
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.CENTER
    prs.save(save_file_path)