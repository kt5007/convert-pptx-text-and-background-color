import pptx
import collections.abc
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import os
from pptx.enum.dml import MSO_THEME_COLOR
import shutil

def convert_text_and_background_color(target_file_path,save_file_path):
    prs = pptx.Presentation(target_file_path)
    slides = prs.slides
    for slide in slides:
        # Change background color to white
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255,255,255)
        for shape in slide.shapes:
            if not shape.has_text_frame: 
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(0, 0, 0)
    prs.save(save_file_path)

if __name__ == '__main__':
    save_directory = './converted_files'
    target_directory = './target_files'
    done_file_directory = './done_files'
    for file_name in os.listdir(target_directory):
        try:
            target_file_path = os.path.join(target_directory,file_name)
            save_file_path = os.path.join(save_directory,file_name)
            # print processing file name
            print(target_file_path)
            # convert
            convert_text_and_background_color(target_file_path,save_file_path)
            # move completed file
            shutil.move(target_file_path, os.path.join(done_file_directory+'/'+file_name))
        except Exception as e:
            with open('exception.txt', 'a') as f:
                print(file_name, file=f)
                print(repr(e), file=f)